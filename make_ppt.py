from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_slide(prs, title_text, bullet_points):
    # Using a bullet point layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title Styling (Smaller and cleaner)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) # White

    # Content Styling
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 200, 200) # Light Gray
        p.space_after = Pt(10)

prs = Presentation()

# --- FIXING THE SIZE TO STANDARD 4:3 ---
# Standard 4:3 size in inches is 10 x 7.5
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# --- SLIDE 1: TITLE ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Autonomous Scientific\nHypothesis Generator"
title.text_frame.paragraphs[0].font.size = Pt(40)
subtitle.text = "LLM-driven Research Synthesis\n\nPresented by: [Your Name]\nGuided by: [Sir's Name]"
subtitle.text_frame.paragraphs[0].font.size = Pt(18)

# --- ADDING CONTENT SLIDES ---
slides_data = [
    ["Introduction", [
        "The 'Knowledge Bottleneck': Millions of papers published annually.",
        "Physical impossibility for humans to track every research gap.",
        "AI assistant bridges existing data and new breakthroughs.",
        "Transform static literature into testable hypotheses."
    ]],
    ["Problem Statement", [
        "Manual literature review is time-intensive (weeks to months).",
        "High probability of human oversight in semantic connections.",
        "Research directions chosen by intuition rather than metrics.",
        "Scalability Crisis: Processing multi-domain datasets in real-time."
    ]],
    ["Proposed System", [
        "AI-Powered Workflow: Leveraging Llama-3.3 and Gemini 2.5.",
        "Automated Gap Discovery: Finding unexplored technical areas.",
        "Objective Metrics: Scientific Impact and Feasibility scores.",
        "Real-time Visualization: Priority Matrix for experimental planning."
    ]],
    ["System Architecture", [
        "Frontend: React.js (Vite) for a high-performance animated SPA.",
        "Backend: FastAPI (Python) for asynchronous AI handling.",
        "Inference Engine: Groq LPU for sub-second speed.",
        "Data Pipeline: PDF -> Text -> LLM Synthesis -> JSON."
    ]],
    ["Technical Stack", [
        "Languages: Python 3.12, JavaScript, Tailwind CSS v4.",
        "Libraries: Framer Motion, Recharts, PyPDF.",
        "AI Models: Meta Llama-3.3-70B via Groq API.",
        "Deployment: Render Cloud Platform."
    ]],
    ["Methodology", [
        "Phase 1: Raw text extraction from multi-page PDFs.",
        "Phase 2: Semantic comprehension of research methodology.",
        "Phase 3: Cross-referencing against scientific knowledge bases.",
        "Phase 4: Algorithmic ranking of discovery potential."
    ]],
    ["System Requirements", [
        "Hardware: Intel i5+, 8GB RAM, SSD Storage.",
        "Software: Windows 10/11, VS Code, Node.js, Python 3.10.",
        "Connectivity: High-speed internet for API inference.",
        "Security: Environment Variable (.env) encryption."
    ]],
    ["Conclusion", [
        "Conclusion: 90% reduction in literature review time.",
        "Future Scope: RAG (Retrieval Augmented Generation) for live data.",
        "Collaboration: Multi-user research workspaces.",
        "Final Goal: Transitioning from raw data to global breakthroughs."
    ]]
]

for s_title, s_points in slides_data:
    add_slide(prs, s_title, s_points)

# --- SLIDE 11: THANK YOU ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Thank You!"
subtitle = slide.placeholders[1]
subtitle.text = "Questions?\n\nContact: [Your Email]"

prs.save('Hypothesis_Generator_Final.pptx')
print("✅ Success! Fixed 4:3 PPT created: Hypothesis_Generator_Final.pptx")